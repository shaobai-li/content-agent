"""Tool: 批量导入文件到指定知识库。

语义等价于 B 版 (build-knowledge-graph) 的 import.py，但作为 Agent Tool
运行在 backend 进程内，无需子进程调用。

参数:
    database_name: 知识库名称（在 databases.json 中注册的名称）
    file_paths:    要导入的文件路径列表
"""
from __future__ import annotations

import hashlib
import json
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from app.agents.tools.base import Tool

# ── 可选依赖 ──────────────────────────────────────────────────────────

try:
    import mammoth

    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from pptx import Presentation

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False


# ── 时间戳 ────────────────────────────────────────────────────────────


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


# ── 文件指纹 ──────────────────────────────────────────────────────────


def _fingerprint_sha256(src: Path) -> str:
    digest = hashlib.sha256()
    with src.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _validate_source_file(src: Path) -> None:
    if not src.exists():
        raise FileNotFoundError(f"输入文件不存在: {src}")
    if not src.is_file():
        raise ValueError(f"输入路径不是文件: {src}")
    try:
        with src.open("rb"):
            pass
    except OSError as exc:
        raise ValueError(f"输入文件不可读: {src}") from exc


# ── ID / 路径 ─────────────────────────────────────────────────────────


def _new_m_id() -> str:
    return str(uuid.uuid4())


def _material_dir(kb_root: Path, m_id: str) -> Path:
    return kb_root / "raw" / f"m_{m_id}"


# ── metadata.json 管理 ────────────────────────────────────────────────


def _metadata_path(kb_root: Path) -> Path:
    return kb_root / "metadata.json"


def _load_metadata(kb_root: Path) -> dict[str, Any]:
    path = _metadata_path(kb_root)
    if not path.exists():
        return {}
    with path.open("r", encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict):
        raise ValueError(f"metadata.json 顶层结构必须是对象: {path}")
    return data


def _save_metadata(kb_root: Path, data: dict[str, Any]) -> None:
    kb_root.mkdir(parents=True, exist_ok=True)
    path = _metadata_path(kb_root)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _is_duplicate_by_fingerprint(metadata: dict[str, Any], sha256_hex: str) -> bool:
    fingerprints = metadata.get("fingerprints")
    if not isinstance(fingerprints, dict):
        return False
    return sha256_hex in fingerprints


# ── record.json 管理 ──────────────────────────────────────────────────


def _build_record(
    m_id: str,
    input_path: Path,
    sha256_hex: str,
    status: str,
    error: str | None = None,
) -> dict[str, Any]:
    record: dict[str, Any] = {
        "m_id": m_id,
        "status": status,
        "created_at": _now_iso(),
        "source": {
            "path": str(input_path),
            "name": input_path.name,
            "size_bytes": input_path.stat().st_size,
            "sha256": sha256_hex,
        },
    }
    if error:
        record["error"] = error
    return record


def _save_record_json(m_dir: Path, record: dict[str, Any]) -> Path:
    m_dir.mkdir(parents=True, exist_ok=True)
    record_path = m_dir / "record.json"
    record_path.write_text(
        json.dumps(record, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return record_path


# ── metadata 变更 ─────────────────────────────────────────────────────


def _append_import_to_metadata(
    metadata: dict[str, Any],
    m_id: str,
    sha256_hex: str,
    input_path: Path,
) -> dict[str, Any]:
    fingerprints = metadata.get("fingerprints")
    if not isinstance(fingerprints, dict):
        fingerprints = {}
        metadata["fingerprints"] = fingerprints

    fingerprints[sha256_hex] = {
        "m_id": m_id,
        "imported_at": _now_iso(),
        "source_name": input_path.name,
    }

    imports = metadata.get("imports")
    if not isinstance(imports, list):
        imports = []
        metadata["imports"] = imports
    imports.append({
        "m_id": m_id,
        "status": "imported",
        "source_name": input_path.name,
        "sha256": sha256_hex,
        "timestamp": _now_iso(),
    })
    metadata["last_import_at"] = _now_iso()
    return metadata


def _append_failure_to_metadata(
    metadata: dict[str, Any],
    m_id: str,
    sha256_hex: str | None,
    error: str,
    input_path: Path | None = None,
) -> dict[str, Any]:
    imports = metadata.get("imports")
    if not isinstance(imports, list):
        imports = []
        metadata["imports"] = imports
    payload: dict[str, Any] = {
        "m_id": m_id,
        "status": "failed",
        "error": error,
        "timestamp": _now_iso(),
    }
    if sha256_hex:
        payload["sha256"] = sha256_hex
    if input_path:
        payload["source_name"] = input_path.name
    imports.append(payload)
    metadata["last_failure_at"] = _now_iso()
    return metadata


# ── view/nodes.json 同步 ─────────────────────────────────────────────


def _sync_import_to_nodes(
    kb_root: Path, m_id: str, input_path: Path, sha256_hex: str
) -> None:
    """将导入的记录同步到 view/nodes.json，使前端能展示该文件。"""
    nodes_path = kb_root / "view" / "nodes.json"
    nodes_path.parent.mkdir(parents=True, exist_ok=True)

    if nodes_path.exists():
        with nodes_path.open("r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = {"kb_id": kb_root.name, "version": 1, "nodes": []}

    if not isinstance(data, dict):
        data = {"kb_id": kb_root.name, "version": 1, "nodes": []}

    nodes = data.get("nodes")
    if not isinstance(nodes, list):
        nodes = []
        data["nodes"] = nodes

    # 幂等：该记录已存在则跳过
    existing_ids = {n.get("record_id") for n in nodes if isinstance(n, dict)}
    if m_id in existing_ids:
        return

    # 确保根文件夹存在
    if not any(isinstance(n, dict) and n.get("id") == "fld_root" for n in nodes):
        now = (
            datetime.now(timezone.utc)
            .replace(microsecond=0)
            .isoformat()
            .replace("+00:00", "Z")
        )
        nodes.insert(0, {
            "id": "fld_root",
            "node_type": "folder",
            "name": "Root",
            "parent_id": None,
            "created_at": now,
            "updated_at": now,
        })

    now = (
        datetime.now(timezone.utc)
        .replace(microsecond=0)
        .isoformat()
        .replace("+00:00", "Z")
    )
    record_node = {
        "id": f"rec_{m_id}",
        "node_type": "record",
        "record_id": m_id,
        "name": input_path.name,
        "file_ext": input_path.suffix,
        "size_bytes": input_path.stat().st_size,
        "sha256": sha256_hex,
        "parent_id": "fld_root",
        "created_at": now,
        "updated_at": now,
    }
    nodes.append(record_node)

    with nodes_path.open("w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


# ── content_paths（B 版特有：记录文档读取路径） ─────────────────────


def _content_paths(src: Path, parsed: dict[str, Any] | None) -> dict[str, str]:
    """返回下一步需要读取的文档路径。

    如果文件被成功解析（PDF/DOCX/PPTX → markdown），返回 parsed.md 路径；
    否则返回原始文件路径。
    """
    if parsed and parsed.get("markdown_path"):
        return {"parsed_path": parsed["markdown_path"]}
    return {"source_path": str(src)}


# ── 解析器（PDF / DOCX / PPTX → markdown） ──────────────────────────


class _DocxParser:
    def parse(self, src_path: Path, output_dir: Path) -> dict[str, str]:
        if not _HAS_DOCX:
            raise ImportError("DOCX 解析需要 mammoth 库：pip install mammoth")
        output_dir.mkdir(parents=True, exist_ok=True)
        result = mammoth.convert_to_markdown(str(src_path))
        markdown_path = output_dir / "parsed.md"
        markdown_path.write_text(result.value, encoding="utf-8")
        return {"markdown_path": str(markdown_path)}


class _PptxParser:
    def parse(self, src_path: Path, output_dir: Path) -> dict[str, str]:
        if not _HAS_PPTX:
            raise ImportError("PPTX 解析需要 python-pptx 库：pip install python-pptx")
        output_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation(str(src_path))
        lines: list[str] = []
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            lines.append(f"## Slide {slide_idx}")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    lines.append(text.strip())
            lines.append("")
        markdown_path = output_dir / "parsed.md"
        markdown_path.write_text("\n".join(lines).strip() + "\n", encoding="utf-8")
        return {"markdown_path": str(markdown_path)}


_PARSER_BY_SUFFIX: dict[str, type] = {
    ".docx": _DocxParser,
    ".pptx": _PptxParser,
}


def _get_parser_for_path(path: Path):
    """根据文件扩展名返回解析器实例，不支持的格式返回 None。"""
    cls = _PARSER_BY_SUFFIX.get(path.suffix.lower())
    return cls() if cls is not None else None


# ── 单文件导入（等价于 B 版 import.py 的 run_import） ────────────────


def _import_single_file(src_path: Path, kb_root: Path) -> dict[str, Any]:
    """导入单个文件到知识库，返回 B 版 import.py 风格的 result dict。"""
    src = src_path.expanduser().resolve()
    kb = kb_root.expanduser().resolve()

    _validate_source_file(src)
    sha256_hex = _fingerprint_sha256(src)

    metadata = _load_metadata(kb)
    if _is_duplicate_by_fingerprint(metadata, sha256_hex):
        return {
            "ok": True,
            "status": "duplicate",
            "reason": "file_already_imported",
            "sha256": sha256_hex,
            "input_file": str(src),
            "file_name": src.name,
            "kb_root": str(kb),
        }

    m_id = _new_m_id()
    m_dir = _material_dir(kb, m_id)

    try:
        parser = _get_parser_for_path(src)
        parsed: dict[str, Any] | None = None

        if parser:
            parsing_record = _build_record(
                m_id=m_id,
                input_path=src,
                sha256_hex=sha256_hex,
                status="parsing",
            )
            _save_record_json(m_dir, parsing_record)
            parsed = parser.parse(src, m_dir)

        record = _build_record(
            m_id=m_id,
            input_path=src,
            sha256_hex=sha256_hex,
            status="imported",
        )
        if parsed:
            record["parsed"] = parsed
        record.update(_content_paths(src, parsed))
        record_path = _save_record_json(m_dir, record)

        updated = _append_import_to_metadata(metadata, m_id, sha256_hex, src)
        _save_metadata(kb, updated)

        # 同步到 view/nodes.json（非关键路径）
        try:
            _sync_import_to_nodes(kb, m_id, src, sha256_hex)
        except Exception:
            pass

        result: dict[str, Any] = {
            "ok": True,
            "status": "imported",
            "m_id": m_id,
            "file_name": src.name,
            "record_path": str(record_path),
            "sha256": sha256_hex,
            "kb_root": str(kb),
            **_content_paths(src, parsed),
        }
        if parsed:
            result["parsed"] = parsed
        return result

    except Exception as exc:
        error_msg = str(exc)
        failure_record = _build_record(
            m_id=m_id,
            input_path=src,
            sha256_hex=sha256_hex,
            status="failed",
            error=error_msg,
        )
        record_path = _save_record_json(m_dir, failure_record)
        failed_meta = _append_failure_to_metadata(
            metadata, m_id, sha256_hex, error_msg, src,
        )
        _save_metadata(kb, failed_meta)
        return {
            "ok": False,
            "status": "failed",
            "m_id": m_id,
            "file_name": src.name,
            "record_path": str(record_path),
            "sha256": sha256_hex,
            "error": error_msg,
            "kb_root": str(kb),
        }


# ═══════════════════════════════════════════════════════════════════════
# Agent Tool
# ═══════════════════════════════════════════════════════════════════════


class ImportKnowledgeTool(Tool):
    """批量导入文件到指定知识库。"""

    def __init__(self, workspace: Path, agent_id: str):
        self._workspace = workspace
        self._agent_id = agent_id

    @property
    def name(self) -> str:
        return "import_knowledge"

    @property
    def description(self) -> str:
        return (
            "将文件批量导入到指定知识库中。支持 DOCX、PPTX、MD、TXT 等格式。"
            "DOCX/PPTX 会自动解析为 Markdown 文本。"
            "注意：PDF 文件请先使用 pdf2md 工具转换为 Markdown 后再导入。"
            "自动基于 SHA256 指纹去重，避免重复导入同一文件。"
            "会同步更新 knowledge_base 目录下的 metadata.json 和 view/nodes.json"
            "供前端知识库面板展示。"
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "database_name": {
                    "type": "string",
                    "description": (
                        "知识库名称，即 databases.json 中注册的名称。"
                        "可通过 list_knowledge_bases 工具查询可用知识库列表。"
                    ),
                },
                "file_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "要导入的文件路径列表。支持绝对路径，"
                        "也支持相对 workspace 目录的相对路径。"
                    ),
                },
            },
            "required": ["database_name", "file_paths"],
        }

    @property
    def read_only(self) -> bool:
        return False

    async def execute(
        self,
        database_name: str,
        file_paths: list[str],
        **kwargs: Any,
    ) -> str:
        # 延迟导入，避免应用启动时的循环依赖
        from app.core.config import get_agent_local_data_dir
        from app.service.knowledge_base_registry_service import list_knowledge_bases

        # ── 1. 查找知识库 ──────────────────────────────────────────
        databases = list_knowledge_bases(self._agent_id)
        kb = next(
            (db for db in databases if db["name"] == database_name),
            None,
        )
        if kb is None:
            available = [db["name"] for db in databases]
            raise ValueError(
                f"知识库「{database_name}」未找到。"
                f"可用知识库：{available}"
            )

        kb_id = kb["id"]
        kb_root = get_agent_local_data_dir(self._agent_id) / kb_id

        # ── 2. 逐个导入文件 ──────────────────────────────────────
        results: list[dict[str, Any]] = []
        for file_path in file_paths:
            src = Path(file_path)
            if not src.is_absolute():
                src = (self._workspace / src).resolve()
            result = _import_single_file(src, kb_root)
            results.append(result)

        # ── 3. 汇总 ──────────────────────────────────────────────
        imported = sum(1 for r in results if r.get("status") == "imported")
        duplicates = sum(1 for r in results if r.get("status") == "duplicate")
        failed = sum(1 for r in results if r.get("status") == "failed")

        summary_parts: list[str] = [
            f"导入知识库「{database_name}」完成，共处理 {len(results)} 个文件",
        ]
        if imported:
            summary_parts.append(f"  ✅ 成功导入：{imported}")
        if duplicates:
            summary_parts.append(f"  ⏭️  重复跳过：{duplicates}")
        if failed:
            summary_parts.append(f"  ❌ 导入失败：{failed}")
        summary_parts.append("")

        for r in results:
            fname = r.get("file_name", "?")
            st = r.get("status", "?")
            if st == "imported":
                cp = r.get("parsed_path") or r.get("source_path", "")
                summary_parts.append(f"  ✅ {fname} → {cp}")
            elif st == "duplicate":
                summary_parts.append(f"  ⏭️  {fname}（重复，已跳过）")
            elif st == "failed":
                summary_parts.append(f"  ❌ {fname}：{r.get('error', '未知错误')}")

        return json.dumps(
            {
                "ok": all(r.get("ok") for r in results),
                "database_name": database_name,
                "kb_id": kb_id,
                "summary": {
                    "total": len(results),
                    "imported": imported,
                    "duplicate": duplicates,
                    "failed": failed,
                },
                "results": results,
                "text": "\n".join(summary_parts),
            },
            ensure_ascii=False,
        )
