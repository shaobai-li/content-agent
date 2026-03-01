"""PPTX to Markdown parser."""
import re
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PptxParser:
    """PPTX文档解析器"""
    
    def __init__(self, position_tolerance: float = 10.0):
        """
        初始化PPTX解析器
        
        Args:
            position_tolerance: 形状定位的容差（点）
        """
        self.position_tolerance = position_tolerance
        self._image_counter = 0
        self._media_dir = None
    
    async def parse(self, file_path: Path, output_dir: Path) -> Path:
        """
        解析PPTX文档为Markdown
        
        Args:
            file_path: PPTX文件路径
            output_dir: Markdown输出目录
            
        Returns:
            生成的Markdown文件路径
        """
        if not file_path.exists():
            raise FileNotFoundError(f"Input file not found: {file_path}")
        
        output_dir.mkdir(parents=True, exist_ok=True)
        
        self._media_dir = output_dir / "media"
        self._media_dir.mkdir(parents=True, exist_ok=True)
        self._image_counter = 0
        
        presentation = Presentation(str(file_path))
        
        md_lines = [
            f"# {file_path.stem}\n\n",
            f"Total slides: {len(presentation.slides)}\n\n"
        ]
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            md_lines.append(f"\n---\n\n## Slide {slide_num}\n\n")
            slide_content = self._extract_slide_content(slide, slide_num)
            md_lines.extend(slide_content)
        
        markdown_content = "".join(md_lines)
        
        md_filename = file_path.stem + ".md"
        md_path = output_dir / md_filename
        md_path.write_text(markdown_content.strip(), encoding="utf-8")
        
        return md_path
    
    def _extract_slide_content(self, slide, slide_num: int) -> List[str]:
        """
        从单个幻灯片提取内容
        
        Args:
            slide: PPTX幻灯片对象
            slide_num: 幻灯片编号
            
        Returns:
            Markdown行列表
        """
        lines = []
        
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"### {slide.shapes.title.text.strip()}\n\n")
        
        shapes = self._collect_and_sort_shapes(slide)
        
        for shape in shapes:
            if shape.has_text_frame:
                lines.extend(self._extract_text_content(shape))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                lines.append(self._extract_image(shape, slide_num))
            elif shape.has_table:
                lines.extend(self._extract_table(shape))
        
        return lines
    
    def _collect_and_sort_shapes(self, slide) -> List:
        """
        收集并按位置排序形状
        
        Args:
            slide: PPTX幻灯片对象
            
        Returns:
            已排序的形状列表
        """
        meaningful_shapes = []
        
        for shape in slide.shapes:
            if (shape.has_text_frame or 
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE or 
                shape.has_table):
                
                top = shape.top.pt if hasattr(shape, "top") else 0
                left = shape.left.pt if hasattr(shape, "left") else 0
                meaningful_shapes.append((top, left, shape))
        
        meaningful_shapes.sort(
            key=lambda x: (x[0] // self.position_tolerance, x[1])
        )
        
        return [shape for _, _, shape in meaningful_shapes]
    
    def _extract_text_content(self, shape) -> List[str]:
        """
        从形状提取文本内容
        
        Args:
            shape: PPTX形状对象
            
        Returns:
            Markdown行列表
        """
        lines = []
        
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            
            if not text:
                continue
            
            level = paragraph.level
            
            if level == 0:
                if re.match(r'^\d+\.\s', text) or re.match(r'^[•\-\*]\s', text):
                    lines.append(f"- {text}\n")
                else:
                    lines.append(f"{text}\n\n")
            else:
                indent = "  " * level
                lines.append(f"{indent}- {text}\n")
        
        return lines
    
    def _extract_image(self, shape, slide_num: int) -> str:
        """
        从形状提取并保存图片
        
        Args:
            shape: PPTX形状对象
            slide_num: 幻灯片编号
            
        Returns:
            Markdown图片引用
        """
        image = shape.image
        ext = "jpg" if image.ext in ("jpeg", "jpg") else image.ext.lower()
        filename = f"slide_{slide_num:02d}_img_{self._image_counter:03d}.{ext}"
        
        image_path = self._media_dir / filename
        image_path.write_bytes(image.blob)
        
        self._image_counter += 1
        
        return f"![](media/{filename})\n\n"
    
    def _extract_table(self, shape) -> List[str]:
        """
        从形状提取表格
        
        Args:
            shape: PPTX形状对象
            
        Returns:
            Markdown表格行列表
        """
        lines = ["\n"]
        table = shape.table
        
        rows_data = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text = cell.text_frame.text.replace('\n', ' ').replace('\r', ' ').strip()
                text = text.replace('|', '\\|')
                cells.append(text)
            rows_data.append(cells)
        
        if rows_data:
            lines.append("| " + " | ".join(rows_data[0]) + " |\n")
            lines.append("| " + " | ".join(["---"] * len(rows_data[0])) + " |\n")
            
            for row_data in rows_data[1:]:
                lines.append("| " + " | ".join(row_data) + " |\n")
            
            lines.append("\n")
        
        return lines

