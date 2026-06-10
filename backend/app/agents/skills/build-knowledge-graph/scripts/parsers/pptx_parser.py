from __future__ import annotations

from pathlib import Path

from pptx import Presentation


class PptxParser:
    def parse(self, src_path: Path, output_dir: Path) -> dict[str, str]:
        output_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation(str(src_path))
        lines: list[str] = []
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            lines.append(f"## Slide {slide_idx}")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    lines.append(text.strip())
            lines.append("")
        markdown_path = output_dir / "parsed.md"
        markdown_path.write_text("\n".join(lines).strip() + "\n", encoding="utf-8")
        return {"markdown_path": str(markdown_path)}
